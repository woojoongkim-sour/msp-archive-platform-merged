from __future__ import annotations

import os
import re
import json
from dataclasses import dataclass
from typing import List, Optional, Tuple

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None  # type: ignore

try:
    from docx import Document as Docx  # python-docx
except Exception:
    Docx = None  # type: ignore

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None  # type: ignore

try:
    import openpyxl  # openpyxl
except Exception:
    openpyxl = None  # type: ignore

try:
    import xlrd  # xlrd for .xls
except Exception:
    xlrd = None  # type: ignore


@dataclass
class ParsedSection:
    title: Optional[str]
    content: str


def _split_paragraphs(text: str) -> List[str]:
    # Simple paragraph splitter: split on blank lines
    parts = re.split(r"\n\n+", text)
    return [p.strip() for p in parts if p and p.strip()]


def parse_document(file_path: str) -> Tuple[List[ParsedSection], bool]:
    """Extract text sections from a document, with optional section titles.

    Returns (sections, meta_only).
    - sections: list of ParsedSection(title, content)
    - meta_only: True if the file could not be parsed due to encryption/protection or error
    """
    ext = os.path.splitext(file_path)[1].lower()
    sections: List[ParsedSection] = []
    meta_only = False

    try:
        if ext == ".pdf":
            if fitz is None:
                raise ImportError("PyMuPDF (fitz) not installed")
            doc = fitz.open(file_path)
            current_title: Optional[str] = None
            current_content_parts: List[str] = []
            for page in doc:
                blocks = page.get_text("dict").get("blocks", [])
                for b in blocks:
                    if b.get("type") != 0:
                        continue
                    for line in b.get("lines", []):
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            if not text:
                                continue
                            # Heuristic: bold + large font as heading
                            font_size = span.get("size", 0)
                            bold = span.get("bold", False) or ("bold" in span.get("font", "").lower())
                            if bold and font_size >= 14:
                                # flush previous
                                if current_content_parts:
                                    sections.append(ParsedSection(title=current_title, content=" ".join(current_content_parts).strip()))
                                    current_content_parts = []
                                current_title = text
                            else:
                                current_content_parts.append(text)
            if current_content_parts:
                sections.append(ParsedSection(title=current_title, content=" ".join(current_content_parts).strip()))
            doc.close()
        elif ext == ".docx":
            if Docx is None:
                raise ImportError("python-docx not installed")
            doc = Docx(file_path)
            current_title: Optional[str] = None
            content_parts: List[str] = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if para.style and "Heading" in para.style.name:
                    if content_parts:
                        sections.append(ParsedSection(title=current_title, content=" ".join(content_parts).strip()))
                        content_parts = []
                    current_title = text
                else:
                    content_parts.append(text)
            if content_parts:
                sections.append(ParsedSection(title=current_title, content=" ".join(content_parts).strip()))
        elif ext in {".ppt", ".pptx"}:
            if Presentation is None:
                raise ImportError("python-pptx not installed")
            prs = Presentation(file_path)
            for slide in prs.slides:
                title = None
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title = slide.shapes.title.text.strip()
                text_parts: List[str] = []
                for shp in slide.shapes:
                    if hasattr(shp, "text_frame") and shp.text:
                        # Collect all text from shapes with text
                        for para in shp.text_frame.paragraphs:
                            text_parts.append(para.text.strip())
                content = " ".join([t for t in text_parts if t]).strip()
                if title or content:
                    sections.append(ParsedSection(title=title, content=content))
        elif ext == ".xlsx":
            if openpyxl is None:
                raise ImportError("openpyxl not installed")
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            for sheet in wb.worksheets:
                title = sheet.title
                lines: List[str] = []
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join([str(v) if v is not None else "" for v in row])
                    lines.append(line)
                content = "\n".join(lines).strip()
                if content:
                    sections.append(ParsedSection(title=title, content=content))
        elif ext == ".xls":
            if xlrd is None:
                raise ImportError("xlrd not installed")
            wb = xlrd.open_workbook(file_path)
            for sheet in wb.sheets():
                title = sheet.name
                lines: List[str] = []
                for r in range(sheet.nrows):
                    row_vals = []
                    for c in range(sheet.ncols):
                        val = sheet.cell(r, c).value
                        row_vals.append(str(val) if val is not None else "")
                    lines.append("\t".join(row_vals))
                content = "\n".join(lines).strip()
                if content:
                    sections.append(ParsedSection(title=title, content=content))
        elif ext in {".txt", ".md", ".csv"}:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
            if text:
                sections.append(ParsedSection(title=None, content=text))
        elif ext == ".hwp":
            # HWP 5.0 — pyhwp (hwp5txt CLI) 로 텍스트 추출
            try:
                import subprocess
                result = subprocess.run(
                    ["hwp5txt", file_path],
                    capture_output=True, text=True, timeout=120,
                )
                if result.returncode == 0:
                    text = (result.stdout or "").strip()
                    if text:
                        sections.append(ParsedSection(title=None, content=text))
                    else:
                        meta_only = True
                else:
                    err = (result.stderr or "").lower()
                    meta_only = True
            except FileNotFoundError:
                meta_only = True  # hwp5txt 미설치
            except subprocess.TimeoutExpired:
                meta_only = True
        else:
            # Fallback: treat as plain text if possible
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read().strip()
                if text:
                    sections.append(ParsedSection(title=None, content=text))
            except Exception:
                meta_only = True
    except Exception:
        meta_only = True

    return sections, meta_only
